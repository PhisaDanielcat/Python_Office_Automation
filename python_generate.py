from docx import Document
from docx.oxml.ns import qn
from copy import deepcopy

from openpyxl import load_workbook

from pptx import Presentation
from pptx.util import Pt,Cm
from pptx.dml.color import RGBColor

import urllib
import time
from tqdm import tqdm

class student:
    def __init__(self,name,sex,home,grade,major,points,awards,adv,disadv,time_apply,time_pos,time_develop,classnumber,pic_url):
        self.name = str(name)
        self.sex = str(sex)
        self.home = str(home)
        self.grade = str(grade)
        self.major = str(major)
        self.points = str(points)
        self.awards = str(awards)
        self.adv = str(adv)
        self.disadv = str(disadv)
        self.time_apply = str(time_apply)
        self.time_pos = str(time_pos)
        self.time_develop = str(time_develop)
        self.classnumber = str(classnumber)
        self.pic_url = str(pic_url)

chinese_year = "2023年"
chinese_date = "3月15日"
student_num = 6
excel_file_path = "origin_file/test6.xlsx"
introducer_file_path = "origin_file/入党介绍人表.xlsx"
word_file_name = "一支部" + chinese_date + "支部委员会意见"
ppt_file_name ="一支部" + chinese_date + "发展对象公示海报"
introducer_file_name = "一支部" + chinese_date + "入党介绍人意见"
general_assembly_meeting_PPT_name = "一支部" + chinese_date + "党员大会PPT"
vote_file_name = "一支部" + chinese_date + "党员大会选票"
presentation_PPT_name = "物院发展和转正汇报-本科一支部" + chinese_year + chinese_date

############### excel operation ##################
workbook = load_workbook(filename = excel_file_path)
sheet = workbook["Sheet1"]

student_list=[]
for i in range(student_num):
    index_list =['G','H','I','J','K','L','M','N','O','P','Q','R','S','T']
    for j in range(len(index_list)):
        index_list[j]+=str(i+2)

    student_object = student(sheet[index_list[0]].value, \
                            sheet[index_list[1]].value, \
                            sheet[index_list[2]].value, \
                            sheet[index_list[3]].value, \
                            sheet[index_list[4]].value, \
                            sheet[index_list[5]].value, \
                            sheet[index_list[6]].value, \
                            sheet[index_list[7]].value, \
                            sheet[index_list[8]].value, \
                            sheet[index_list[9]].value, \
                            sheet[index_list[10]].value, \
                            sheet[index_list[11]].value, \
                            sheet[index_list[12]].value, \
                            sheet[index_list[13]].value)
    student_list.append(student_object)

introducer_workbook = load_workbook(filename=introducer_file_path)
sheet = introducer_workbook["Sheet1"]
for i in range(student_num):
    index_list = ['A','B','C']
    for j in range(len(index_list)):
        index_list[j]+=str(i+2)
    for k in range(student_num):
        if(student_list[k].name == sheet[index_list[0]].value):
            student_list[k].introducor1 = sheet[index_list[1]].value
            student_list[k].introducor2 = sheet[index_list[2]].value

print("downloading photos.")
now_time = time.time()


for i in tqdm(range(student_num)):
    urlStr = student_list[i].pic_url
    pic_file_name = './pictures/' + student_list[i].name + '.png'
    urllib.request.urlretrieve(urlStr, filename = pic_file_name)


print("photos already downloaded.")
print("excel extraction done.")
print("All files generated start.")
time.sleep(1)

############### excel operation done##################

############### word operation #######################
doc = Document("temp/temp.docx")
doc.styles['Normal'].font.name = 'Times New Roman'
doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')

student_title_list=[]
student1_group_eva_list=[]
for i in range(student_num):
    student_title=student_list[i].name + '：'
    student_title_list.append(student_title)
    student1_group_eva = '物理科学与技术学院本科生第一党支部于',chinese_date,'召开支部委员会讨论研究了关于'+ student_list[i].name + '同志的发展事宜。' \
                     + student_list[i].name + '同志'+ student_list[i].time_apply + '提交入党申请书，于'+ student_list[i].time_pos + '成为入党积极分子，第' \
                     + student_list[i].classnumber + '发展对象培训班结业，并在'+ student_list[i].time_develop + '成为发展对象。'+ student_list[i].name + '曾获'+ student_list[i].awards + '等荣誉。优点是' \
                     + student_list[i].adv + '，但存在'+ student_list[i].disadv + '的缺点。总体来看，经支部委员会研究，认为'+ student_list[i].name + '同志已基本符合预备党员条件，同意如期发展，提交支部党员大会讨论。'
    student1_group_eva_list.append(student1_group_eva)

for i in range(student_num):
    param1=doc.add_paragraph().add_run()
    param1.text = student_title_list[i]
    param1.font.bold = True
    param1.font.underline = True
    doc.add_paragraph(student1_group_eva_list[i])
    doc.add_paragraph('\r')

word_file_name = "一支部" + chinese_date + "支部委员会意见"
word_file_path = "doc\\" + word_file_name + ".docx"
doc.save(word_file_path)
print(word_file_name,".docx \t\thas generate done.",sep='')
############### word operation done ##################

############### ppt operation #######################
ppt = Presentation('temp/post_temp.pptx')
slide = ppt.slides[0]

# 内容
pre_words = "经党组织一年以上培养教育和考察，"+student_list[0].name+"等"+str(student_num)+ \
            "名同志被确定为发展对象，党支部拟于近期\n召开党员大会讨论其发展问题。为进一步增强发展党员工作的透明度，确保新发展党员质量\n，根据相关规定，现将发展对象相关情况公示如下："
title_paragraph = slide.shapes.add_textbox(left=Cm(0.71875), top=Cm(4.5), width=Cm(16.96875), height=Cm(1.59375)).text_frame
title_paragraph.paragraphs[0].text = pre_words
title_paragraph.paragraphs[0].font.size = Pt(12)
title_paragraph.paragraphs[0].font.name = '微软雅黑'
title_paragraph.paragraphs[0].font.bold = False
title_paragraph.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

for page in range(student_num//3 + 1):
    slide = ppt.slides[page]
    if page==(student_num//3):
        student_each_page = student_num - (student_num//3)*3
    else:
        student_each_page = 3
    for i in range(page*3,page*3 + student_each_page):
        student_words = student_list[i].name + '，'+ student_list[i].sex + '，' +  student_list[i].home + '人，平均绩点' + student_list[i].grade + '，' + student_list[i].awards + '，' \
                        + student_list[i].time_apply + '申请入党，'+ student_list[i].time_pos + '成为入党积极分子并参加第，'+ student_list[i].classnumber + '期党校培训班学习，'  \
                        + student_list[i].time_develop + '成为入党积极分子。'
        new_paragraph1 = slide.shapes.add_textbox(left=Cm(0.71875), top=Cm(7.5 + i*6 - page*18), width=Cm(16.96875), height=Cm(1.59375)).text_frame
        new_paragraph1.paragraphs[0].text = student_words
        new_paragraph1.paragraphs[0].font.size = Pt(12)
        new_paragraph1.paragraphs[0].font.name = '微软雅黑'
        new_paragraph1.paragraphs[0].font.bold = False
        new_paragraph1.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

        pic_path = 'pictures/' + student_list[i].name + '.png'

        slide.shapes.add_picture(pic_path, Cm(13), Cm(7.5 + i*6 - page*18), Cm(4), Cm(5))

ppt_file_name_full = 'doc/' + ppt_file_name + '.pptx'
ppt.save(ppt_file_name_full)
print(ppt_file_name,".ppt \t\thas generate done",sep='')
############### ppt operation #######################

############### inrtoducor opinion########################
doc = Document("temp/temp.docx")
doc.styles['Normal'].font.name = 'Times New Roman'
doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')

introducer_full_name_list = []
for i in range(student_num):
    introducer_full_name_list.append(student_list[i].introducor1)
    introducer_full_name_list.append(student_list[i].introducor2)
introducer_full_name_not_repeat=[]
for i in introducer_full_name_list:
    if i not in introducer_full_name_not_repeat:
        introducer_full_name_not_repeat.append(i)

introducer_title_words = '到场发表意见的党员有：'
for i in range(len(introducer_full_name_not_repeat)):
    introducer_title_words+=(str(introducer_full_name_not_repeat[i])+",")

param1=doc.add_paragraph().add_run()
param1.text = introducer_title_words
param1.font.bold = True
param1.font.underline = True
doc.add_paragraph('\r')

for i in range(student_num):
    intro_selfname=doc.add_paragraph().add_run()
    intro_selfname.text = str(i+1)+'：'+ student_list[i].name
    intro_selfname.font.bold = True

    introducer_1_words = '入党介绍人一' + str(student_list[i].introducor1) + "：" + student_list[i].name+'同志的优点是：'+ \
                         student_list[i].adv+'，但是缺点是：'+student_list[i].disadv+'。综上我推荐'+ student_list[i].name+"同志加入中国共产党。"
    doc.add_paragraph(introducer_1_words)
    introducer_2_words = '入党介绍人一' + str(student_list[i].introducor2) + "：" + student_list[i].name+'同志的优点是：'+ \
                         student_list[i].adv+'，但是缺点是：'+student_list[i].disadv+'。综上我推荐'+ student_list[i].name+"同志加入中国共产党。"
    doc.add_paragraph(introducer_2_words)
    doc.add_paragraph('\r')

introducor_file_name_full = 'doc/'+introducer_file_name+'.docx'
doc.save(introducor_file_name_full)
print(introducer_file_name+".docx \t\thas generate done")

############### inrtoducor opinion done ########################

############### General Assembly PPT ########################
ppt = Presentation('temp/General Assembly PPT.pptx')
slide = ppt.slides[3]

meeting_content = '讨论关于'
for i in range(student_num):
    meeting_content+=student_list[i].name
    if not (i == student_num-1):
        meeting_content += '，'
meeting_content += ('共'+str(student_num)+'位同志的入党事宜')

title_paragraph = slide.shapes.add_textbox(left=Cm(1.1875), top=Cm(7.5625), width=Cm(21.84375), height=Cm(8.375)).text_frame
title_paragraph.paragraphs[0].text = meeting_content
title_paragraph.paragraphs[0].font.size = Pt(32)
title_paragraph.paragraphs[0].font.name = '宋体'
title_paragraph.paragraphs[0].font.bold = True
title_paragraph.paragraphs[0].font.color.rgb = RGBColor(0,0,0)


general_assembly_meeting_PPT_name_full = 'doc/' + general_assembly_meeting_PPT_name + '.pptx'
ppt.save(general_assembly_meeting_PPT_name_full)
print(general_assembly_meeting_PPT_name,".ppt \t\t\thas generate done",sep='')
############### General Assembly PPT generated########################

############### Vote generate ########################################
doc = Document("temp/vote_temp.docx")
doc.styles['Normal'].font.name = 'Times New Roman'
doc.styles['Normal']._element.rPr.rFonts.set(qn('w:eastAsia'), u'宋体')

# 第一个修改
vote_title = doc.paragraphs[0]
vote_title.text = "关于 ",student_list[0].name," 同志入党的表决选票"
vote_title_runs = vote_title.runs[0]
vote_title_runs.font.bold = True
vote_title_runs.font.size = Pt(16)

vote_title_word1 = doc.paragraphs[1]
vote_title_word1.text = "单位：物理科学与技术学院本科生第一党支部"
vote_title_word1_runs = vote_title_word1.runs[0]
vote_title_word1_runs.font.size = Pt(12)

vote_title_word2 = doc.paragraphs[2]
vote_title_word2.text = "投票时间：",chinese_year,chinese_date
vote_title_word2_runs = vote_title_word2.runs[0]
vote_title_word2_runs.font.size = Pt(12)

vote_table = doc.tables[0]
vote_table_name = vote_table.cell(1, 0).paragraphs[0]
vote_table_name.text = student_list[0].name

# 后面的复制

for i in range(1,student_num):
    new_vote_title = deepcopy(vote_title)
    new_vote_title.text = "关于",student_list[i].name,"同志入党的表决选票"
    new_vote_title_runs = new_vote_title.runs[0]
    new_vote_title_runs.font.bold = True
    new_vote_title_runs.font.size = Pt(16)
    paragraph = doc.add_paragraph()
    paragraph._p.addnext(new_vote_title._element)

    new_vote_title_word1 = deepcopy(vote_title_word1)
    paragraph = doc.add_paragraph()
    paragraph._p.addnext(new_vote_title_word1._element)

    new_vote_title_word2 = deepcopy(vote_title_word2)
    paragraph = doc.add_paragraph()
    paragraph._p.addnext(new_vote_title_word2._element)

    vote_table = doc.tables[0]
    new_table = deepcopy(vote_table)
    new_vote_table_name = new_table.cell(1, 0).paragraphs[0]
    new_vote_table_name.text = student_list[i].name
    paragraph = doc.add_paragraph()
    paragraph._p.addnext(new_table._element)

    for j in range(8):
        doc.add_paragraph("\n")
vote_file_name_full = "doc/"+vote_file_name+".docx"
doc.save(vote_file_name_full)
print(vote_file_name,".docx \t\thas generate done")
#################Vote generate done#######################

########## Generate Presentation PPT #####################
ppt = Presentation('temp/presetation ppt.pptx')

slide = ppt.slides[1]
title_paragraph = slide.shapes.add_textbox(left=Cm(301/20), top=Cm(215/20), width=Cm(135/20), height=Cm(36/20)).text_frame
title_paragraph.paragraphs[0].text = "发展党员"+str(student_num)+"名"
title_paragraph.paragraphs[0].font.size = Pt(28)
title_paragraph.paragraphs[0].font.name = '宋体'
title_paragraph.paragraphs[0].font.bold = True
title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192,0,0)

slide = ppt.slides[2]
title_paragraph = slide.shapes.add_textbox(left=Cm(33/20), top=Cm(132/20), width=Cm(206/20), height=Cm(39/20)).text_frame
title_paragraph.paragraphs[0].text = "发展预备党员（"+str(student_num)+"名）"
title_paragraph.paragraphs[0].font.size = Pt(28)
title_paragraph.paragraphs[0].font.name = '宋体'
title_paragraph.paragraphs[0].font.bold = True
title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192,0,0)
student_list_str=''
for i in range(student_num):
    if(i==student_num):
        student_list_str+=student_list[i].name
    else:
        student_list_str += (student_list[i].name + "，")
title_paragraph = slide.shapes.add_textbox(left=Cm(33/20), top=Cm(187/20), width=Cm(206/20), height=Cm(39/20)).text_frame
title_paragraph.paragraphs[0].text = student_list_str
title_paragraph.paragraphs[0].font.size = Pt(24)
title_paragraph.paragraphs[0].font.name = '宋体'
title_paragraph.paragraphs[0].font.bold = False
title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192,0,0)

for i in range(student_num):
    slide = ppt.slides[i+3]
    pic_path = 'pictures/' + student_list[i].name + '.png'
    slide.shapes.add_picture(pic_path, Cm(29/20), Cm(69/20), Cm(163/20), Cm(231/20))

    ppt_student_info = student_list[i].name + "，" + student_list[i].sex + "，平均绩点" + student_list[i].grade \
    + "，获得过" +student_list[i].awards + "等奖项"
    title_paragraph = slide.shapes.add_textbox(left=Cm(190 / 20), top=Cm(71 / 20), width=Cm(320 / 20),
                                               height=Cm(34 / 20)).text_frame
    title_paragraph.paragraphs[0].text = ppt_student_info
    title_paragraph.paragraphs[0].font.size = Pt(18)
    title_paragraph.paragraphs[0].font.name = '宋体'
    title_paragraph.paragraphs[0].font.bold = True
    title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

    ppt_student_adv = "优点" + student_list[i].adv
    title_paragraph = slide.shapes.add_textbox(left=Cm(190 / 20), top=Cm(114 / 20), width=Cm(320 / 20),
                                               height=Cm(34 / 20)).text_frame
    title_paragraph.paragraphs[0].text = ppt_student_adv
    title_paragraph.paragraphs[0].font.size = Pt(18)
    title_paragraph.paragraphs[0].font.name = '宋体'
    title_paragraph.paragraphs[0].font.bold = True
    title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

    ppt_student_disadv = "缺点" + student_list[i].disadv
    title_paragraph = slide.shapes.add_textbox(left=Cm(190 / 20), top=Cm(125 / 20), width=Cm(320 / 20),
                                               height=Cm(34 / 20)).text_frame
    title_paragraph.paragraphs[0].text = ppt_student_disadv
    title_paragraph.paragraphs[0].font.size = Pt(18)
    title_paragraph.paragraphs[0].font.name = '宋体'
    title_paragraph.paragraphs[0].font.bold = True
    title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

    ppt_student_time = student_list[i].time_apply + "申请入党" + "\n" \
                    +  student_list[i].time_pos + "确定为入党积极分子" + "\n" \
                    +  "参加第"+student_list[i].classnumber + "期党校培训班学习" + "\n" \
                    +  student_list[i].time_develop + "确定为发展对象"
    title_paragraph = slide.shapes.add_textbox(left=Cm(190 / 20), top=Cm(163 / 20), width=Cm(320 / 20),
                                               height=Cm(34 / 20)).text_frame
    title_paragraph.paragraphs[0].text = ppt_student_time
    title_paragraph.paragraphs[0].font.size = Pt(18)
    title_paragraph.paragraphs[0].font.name = '宋体'
    title_paragraph.paragraphs[0].font.bold = True
    title_paragraph.paragraphs[0].font.color.rgb = RGBColor(192, 0, 0)

presentation_PPT_name_full = "doc/" + presentation_PPT_name + ".pptx"
ppt.save(presentation_PPT_name_full)
print(presentation_PPT_name,".pptx \t\thas generate done")
########## Generate Presentation PPT done#################


end_time = time.time()
using_time = end_time - now_time
print("All files have generated")
print("Time cost is %.2f"%using_time,"s")